import os
import zipfile
from docx import Document
from pptx import Presentation

BASE = "MRAG_HC_Submission"
os.makedirs(BASE, exist_ok=True)

# -----------------------------
# 1. CREATE IEEE PAPER (DOCX)
# -----------------------------
doc = Document()
doc.add_heading("MRAG-HC System", 0)

doc.add_heading("Abstract", 1)
doc.add_paragraph(
"Multilingual Retrieval-Augmented Generation system with hallucination control "
"for domain-specific knowledge retrieval."
)

doc.add_heading("Introduction", 1)
doc.add_paragraph("LLMs suffer from hallucination issues. This system mitigates them using RAG + verification.")

doc.add_heading("Methodology", 1)
doc.add_paragraph(
"Pipeline includes ingestion, chunking, embedding, retrieval, LLM generation, and verification layer."
)

doc.add_heading("Hallucination Control", 1)
doc.add_paragraph(
"Fact checking + confidence scoring ensures reliability of generated responses."
)

doc.add_heading("Evaluation", 1)
doc.add_paragraph("Metrics include accuracy, faithfulness, and hallucination rate.")

paper_path = os.path.join(BASE, "IEEE_Paper.docx")
doc.save(paper_path)

# -----------------------------
# 2. CREATE PPT
# -----------------------------
ppt = Presentation()

def slide(title, text):
    s = ppt.slides.add_slide(ppt.slide_layouts[1])
    s.shapes.title.text = title
    s.placeholders[1].text = text

slide("MRAG-HC System", "Major Project - M.Tech CSE")
slide("Problem", "LLMs hallucinate incorrect information")
slide("Objective", "Build grounded multilingual AI system")
slide("Architecture", "RAG pipeline with verification layer")
slide("Modules", "Loader, Embedder, Retriever, Generator")
slide("Hallucination Control", "Fact checking + confidence scoring")
slide("Evaluation", "Accuracy, Faithfulness, Hallucination Rate")
slide("Tech Stack", "Python, FAISS, FastAPI, Streamlit")
slide("Conclusion", "Reliable AI knowledge assistant")

ppt_path = os.path.join(BASE, "Presentation.pptx")
ppt.save(ppt_path)

# -----------------------------
# 3. CREATE ZIP PACKAGE
# -----------------------------
zip_path = "MRAG_HC_FINAL_SUBMISSION.zip"

with zipfile.ZipFile(zip_path, "w") as z:
    z.write(paper_path)
    z.write(ppt_path)

print("DONE ✔ Submission Pack Created:", zip_path)
